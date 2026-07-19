# -*- coding: utf-8 -*-
"""產生 PPA AI Hackathon 完整簡報（深色 PPA 風格）。"""
import sys, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

OUT = r"C:\Users\Fung N2\AWS\PPA_AI_Hackathon_簡報.pptx"

# 顏色
NAVY   = RGBColor(0x0A, 0x17, 0x30)
NAVY2  = RGBColor(0x0E, 0x20, 0x44)
CARD   = RGBColor(0x13, 0x29, 0x4B)
ORANGE = RGBColor(0xF7, 0xA8, 0x1B)
CYAN   = RGBColor(0x25, 0xD0, 0xEE)
GREEN  = RGBColor(0x3D, 0xDC, 0x97)
WHITE  = RGBColor(0xEC, 0xF2, 0xFB)
GREY   = RGBColor(0x9F, 0xB2, 0xCE)
RED    = RGBColor(0xF2, 0x6D, 0x6D)
FONT   = "Microsoft JhengHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def _set_cjk(run, font=FONT):
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=None, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold) 或 str"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    if isinstance(lines, str):
        lines = [(lines, 18, WHITE, False)]
    first = True
    for item in lines:
        text, size, color, bold = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        _set_cjk(r)
    return tb

def bullets(slide, x, y, w, h, items, size=16, color=WHITE, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, c, b = it
        else:
            text, c, b = it, color, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "•  " + text
        r.font.size = Pt(size); r.font.color.rgb = c; r.font.bold = b
        _set_cjk(r)
    return tb

def header(slide, kicker, title, num=None):
    rect(slide, 0, 0, SW, Inches(0.12), ORANGE)
    if kicker:
        txt(slide, Inches(0.6), Inches(0.35), Inches(11), Inches(0.4),
            [(kicker, 14, CYAN, True)])
    txt(slide, Inches(0.6), Inches(0.62), Inches(12), Inches(0.9),
        [(title, 30, WHITE, True)])
    if num:
        txt(slide, Inches(12.0), Inches(0.3), Inches(0.9), Inches(0.6),
            [(num, 22, RGBColor(0x2b,0x45,0x74), True)], align=PP_ALIGN.RIGHT)

def footer(slide, idx):
    txt(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
        [("PressPlay AI 智慧學習生態系 · 團隊 fish", 10, RGBColor(0x5a,0x6f,0x93), False)])
    txt(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35),
        [(str(idx), 10, RGBColor(0x5a,0x6f,0x93), False)], align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, color=CARD, line=None):
    return rect(slide, x, y, w, h, color, line=line, radius=True)

_idx = [0]
def new(color=NAVY):
    s = prs.slides.add_slide(BLANK)
    bg(s, color)
    _idx[0] += 1
    return s

# ---------------- 1. 封面 ----------------
s = new(NAVY)
rect(s, 0, 0, SW, Inches(0.16), ORANGE)
rect(s, Inches(0.6), Inches(2.0), Inches(1.0), Inches(1.0), ORANGE, radius=True)
txt(s, Inches(0.6), Inches(2.05), Inches(1.0), Inches(0.9), [("PP", 34, NAVY, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.9), Inches(1.9), Inches(11), Inches(1.4),
    [("PressPlay AI 智慧學習生態系升級", 40, WHITE, True),
     ("從對話到多元應用：主動進化 × 多維賦能的 AI 智慧對話助教", 20, CYAN, False)])
txt(s, Inches(0.6), Inches(4.3), Inches(11), Inches(1.6),
    [("讓被證實有效的 AI 助教（用了影片觀看 +700%），", 18, GREY, False),
     ("從『少數會員的被動客服』進化為『主動、多維、跨場景』的學習引擎。", 18, GREY, False)])
txt(s, Inches(0.6), Inches(6.2), Inches(11), Inches(0.8),
    [("團隊 fish  ｜  AWS Summit Taipei 2026 × PressPlay Academy Hackathon", 16, WHITE, True)])
footer(s, _idx[0])

# ---------------- 2. 命題與核心挑戰 ----------------
s = new()
header(s, "THE CHALLENGE", "命題與核心挑戰", "01")
card(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.1), NAVY2)
txt(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.9),
    [("「從對話到多元應用：打造具備主動進化與多維賦能的 AI 智慧對話助教」", 20, ORANGE, True)],
    anchor=MSO_ANCHOR.MIDDLE)
# 兩張數據卡
card(s, Inches(0.6), Inches(3.0), Inches(5.9), Inches(2.1), CARD)
txt(s, Inches(0.9), Inches(3.2), Inches(5.3), Inches(0.5), [("AI 助教 DAU / MAU", 16, ORANGE, True)])
txt(s, Inches(0.9), Inches(3.7), Inches(5.3), Inches(1.2), [("29%  →  <10%", 40, WHITE, True)])
txt(s, Inches(0.9), Inches(4.7), Inches(5.3), Inches(0.4), [("用的人在流失", 14, GREY, False)])
card(s, Inches(6.8), Inches(3.0), Inches(5.9), Inches(2.1), CARD)
txt(s, Inches(7.1), Inches(3.2), Inches(5.3), Inches(0.5), [("使用 AI 助教的學習深度", 16, ORANGE, True)])
txt(s, Inches(7.1), Inches(3.7), Inches(5.3), Inches(1.3),
    [("平均影片觀看次  +700%", 18, GREEN, True),
     ("平均 Session 數  +84%", 16, WHITE, False),
     ("平均活躍天數  +52%", 16, WHITE, False)])
card(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.1), NAVY2, line=ORANGE)
txt(s, Inches(0.9), Inches(5.45), Inches(11.5), Inches(0.9),
    [("核心挑戰：如何將 AI 的技術與服務，擴大至更多元化的學習場景中？", 18, WHITE, True)],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, _idx[0])

# ---------------- 3. 矛盾核心 ----------------
s = new()
header(s, "THE INSIGHT", "矛盾核心：效果極強，觸及極弱", "02")
txt(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
    [("AI 助教不是沒用 —— 是用的人在流失。真正要解的是「擴散」與「留住」。", 18, CYAN, True)])
data = [
    ("AI 助教滲透率", "11.7%", "被動、會員專屬、等你發問", RED),
    ("知識挑戰滲透率", "42%", "人氣高，但流量停在挑戰頁", ORANGE),
    ("知識挑戰 DAU/MAU", "27%", "沒有回流到 AI 與課程", ORANGE),
]
x = Inches(0.6)
for title, big, sub, col in data:
    card(s, x, Inches(2.5), Inches(3.9), Inches(2.0), CARD)
    txt(s, x + Inches(0.25), Inches(2.65), Inches(3.4), Inches(0.5), [(title, 15, GREY, True)])
    txt(s, x + Inches(0.25), Inches(3.1), Inches(3.4), Inches(0.9), [(big, 36, col, True)])
    txt(s, x + Inches(0.25), Inches(4.0), Inches(3.4), Inches(0.5), [(sub, 12, GREY, False)])
    x += Inches(4.05)
card(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.7), NAVY2, line=ORANGE)
txt(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.4),
    [("⚑ 官方拋出的質疑（本案必須正面回答）", 16, ORANGE, True),
     ("「遊戲化滲透率雖高，但無法確認能讓用戶看更多影片。」", 18, WHITE, True),
     ("→ 我們的中心假設：用機制 + 數據，證明「遊戲化 → 帶動看課 / 拉高滲透率」。", 15, GREY, False)])
footer(s, _idx[0])

# ---------------- 4. 斷點鐵證（實跑資料）----------------
s = new()
header(s, "DATA-DRIVEN", "我們實跑 Dataset 的斷點鐵證（可複現）", "03")
rows = [
    ("39.6%", "的 AI 對話，用戶只問一次就離開", CYAN),
    ("6.68 則", "平均每段對話訊息數（約 3.4 個來回）", WHITE),
    ("13.2% / 20.5%", "行為僅 13.2% 進『學習中』，20.5% 停在『其他』漫遊", ORANGE),
    ("145 類", "課程分類長尾 —— 通用推薦難命中，需 AI 精準分眾", GREEN),
]
y = Inches(1.8)
for big, desc, col in rows:
    card(s, Inches(0.6), y, Inches(12.1), Inches(1.05), CARD)
    txt(s, Inches(0.9), y + Inches(0.12), Inches(3.4), Inches(0.8), [(big, 26, col, True)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.4), y + Inches(0.12), Inches(8.0), Inches(0.8), [(desc, 16, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.16)
txt(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
    [("資料集為 116 位會員 / 159 段對話的取樣，僅用於定位問題，不反推全站指標。", 11, GREY, False)])
footer(s, _idx[0])

# ---------------- 5. 產品戰略：三支柱 ----------------
s = new()
header(s, "STRATEGY", "產品戰略：三支柱，直接對應評分", "04")
pillars = [
    ("① 主動進化", "AI 從『等發問』變『主動接住斷點』：問完 / 答錯 / 看完，主動生成下一步", "創新度 25% + 學習深度", CYAN),
    ("② 多維賦能", "同一份課程 → 挑戰題 / 45秒精華 / 電子書 / Podcast / 雷達診斷", "創新度 25% + 資料洞察 20%", GREEN),
    ("③ 滲透擴散", "AI 限會員 → 45秒軟著陸 + 遊戲化回訪，讓會員更常用、非會員入會", "商業影響力 30%", ORANGE),
]
x = Inches(0.6)
for title, desc, tag, col in pillars:
    card(s, x, Inches(1.9), Inches(3.9), Inches(4.2), CARD, line=col)
    rect(s, x, Inches(1.9), Inches(3.9), Inches(0.12), col)
    txt(s, x + Inches(0.3), Inches(2.2), Inches(3.3), Inches(0.7), [(title, 22, col, True)])
    txt(s, x + Inches(0.3), Inches(3.0), Inches(3.3), Inches(2.0), [(desc, 15, WHITE, False)])
    card(s, x + Inches(0.3), Inches(5.3), Inches(3.3), Inches(0.6), NAVY2)
    txt(s, x + Inches(0.3), Inches(5.35), Inches(3.3), Inches(0.5), [(tag, 12, col, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += Inches(4.05)
footer(s, _idx[0])

# ---------------- 6. 戰略邏輯 flow ----------------
s = new()
header(s, "THE LOOP", "戰略邏輯：努力補償 → 無縫擴散", "05")
steps = ["高頻打開\n(主動診斷)", "答錯/看完\nAI 主動介入", "45秒精華\n軟著陸", "認知懸念\n引發渴望", "續看 / 入會\n擴大滲透池"]
x = Inches(0.6); w = Inches(2.15)
for i, st in enumerate(steps):
    col = ORANGE if i in (2,4) else CARD
    tcol = NAVY if i in (2,4) else WHITE
    card(s, x, Inches(2.6), w, Inches(1.6), col)
    txt(s, x, Inches(2.7), w, Inches(1.4), [(st, 15, tcol, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps)-1:
        txt(s, x + w, Inches(2.9), Inches(0.45), Inches(1.0), [("→", 28, CYAN, True)], align=PP_ALIGN.CENTER)
    x += w + Inches(0.45)
card(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.5), NAVY2, line=ORANGE)
txt(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.2),
    [("關鍵校正：入會不是變現故事，而是擴散手段。", 17, ORANGE, True),
     ("AI 助教是會員專屬 → 每多一位會員 = 多一位 AI 可觸及用戶 = 滲透率分母池被實質放大。", 15, WHITE, False)])
footer(s, _idx[0])

# ---------------- 7~11. 五大模組 ----------------
modules = [
    ("模組 A（主線）", "AI 主動診斷與「下一步」引擎", CYAN,
     ["鎖定問題：39.6% 對話問一次就走、僅 13.2% 進入學習中",
      "串接①提問 ②對話 ⑥軌跡，對話/挑戰結束時主動生成「知識缺口 + 下一步」卡片",
      "主動時機：答錯、看完一段、連續同主題、N 天未回訪",
      "命中官方要的『主動進化』——不再被動等待，接住每個流失斷點"],
     "解決『問完就走』"),
    ("模組 B", "動態懸念剪輯與「45 秒軟著陸」", GREEN,
     ["答錯/卡關時 AI 主動介入，不丟鎖死的付費牆連結",
      "RAG 精準抓對應『45~60 秒免費精華錨點』作試吃墊",
      "最佳時刻懸念截斷（Cliffhanger），引發解鎖渴望",
      "會員一鍵續看；非會員看完引導入會 → 擴大 AI 觸及池"],
     "降低學習摩擦 + 擴散滲透"),
    ("模組 C", "AI 記憶去重與熱門加乘", ORANGE,
     ["呼叫④觀看紀錄，Bedrock 推薦前 Filter 掉已看 >80% 的內容 ID",
      "徹底解決 Beta 版『每次推薦都一樣』的敷衍感",
      "結合③近7天觀看數，命中飆升課程給 1.5 倍學習獎勵",
      "用遊戲化主動把流量導向平台爆款"],
     "解決敷衍感 + 主動導流"),
    ("模組 D", "多模態學習與權限分級", CYAN,
     ["短影音下方『展開看更多』",
      "免費用戶：片段文字",
      "付費會員：AI 轉換『全文本電子書』與『純聲 Podcast』",
      "同一套課程 → 多種學習場景，直接命中『多元化學習場景』"],
     "多元學習場景"),
    ("模組 A 延伸", "知識精靈養成 + 學習戰隊", GREEN,
     ["設計鐵律：遊戲化必須是『看課引擎』，不是平行娛樂",
      "經驗值只能靠『看完課程』或『用 AI 助教』取得（正面回答官方質疑）",
      "用⑤挑戰分類次數養成分眾寵物 PiPi（投資理財→西裝 PP）",
      "2–5 人戰隊，AI 教練派『一起看完 X 課』任務，社交壓力驅動回訪"],
     "把『玩完就走』變『每日回訪』"),
]
for i, (kicker, title, col, items, tagline) in enumerate(modules):
    s = new()
    header(s, kicker, title, "0"+str(6+i) if 6+i < 10 else str(6+i))
    card(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.7), NAVY2, line=col)
    txt(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.6), [("🎯 " + tagline, 17, col, True)], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(3.6), items, size=17, gap=14)
    footer(s, _idx[0])

# ---------------- 12. 資料運用 ----------------
s = new()
header(s, "DATA UTILIZATION", "如何運用官方 6 張資料表", "11")
tbl_rows = [
    ("① AI助教提問", "提問主題、知識缺口聚類", "模組A"),
    ("② AI助教對話", "找『對話結束就流失』斷點、生成下一步", "模組A"),
    ("③ 課程內容(近7/14/30天)", "熱門加乘、45秒錨點來源", "模組B/C"),
    ("④ 會員觀看紀錄", "去識別化行為矩陣、推薦去重", "模組C"),
    ("⑤ 知識挑戰(分類次數)", "寵物經驗值、分眾能力值", "模組A延伸"),
    ("⑥ 接觸AI前後軌跡", "量化行為漏斗、定位漫遊/流失斷點", "模組A+成效"),
]
y = Inches(1.75)
# header row
rect(s, Inches(0.6), y, Inches(4.2), Inches(0.5), NAVY2)
rect(s, Inches(4.85), y, Inches(5.4), Inches(0.5), NAVY2)
rect(s, Inches(10.3), y, Inches(2.4), Inches(0.5), NAVY2)
txt(s, Inches(0.7), y, Inches(4.0), Inches(0.5), [("資料表", 13, CYAN, True)], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(4.95), y, Inches(5.2), Inches(0.5), [("本案用途", 13, CYAN, True)], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(10.4), y, Inches(2.2), Inches(0.5), [("模組", 13, CYAN, True)], anchor=MSO_ANCHOR.MIDDLE)
y += Inches(0.55)
for a, b, c in tbl_rows:
    rect(s, Inches(0.6), y, Inches(4.2), Inches(0.62), CARD)
    rect(s, Inches(4.85), y, Inches(5.4), Inches(0.62), CARD)
    rect(s, Inches(10.3), y, Inches(2.4), Inches(0.62), CARD)
    txt(s, Inches(0.7), y, Inches(4.0), Inches(0.62), [(a, 13, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.95), y, Inches(5.2), Inches(0.62), [(b, 12, GREY, False)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(10.4), y, Inches(2.2), Inches(0.62), [(c, 12, ORANGE, True)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.67)
footer(s, _idx[0])

# ---------------- 13. 成效衡量 ----------------
s = new()
header(s, "SUCCESS METRICS", "成效衡量（對齊商業影響力 30%）", "12")
card(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(2.3), CARD, line=CYAN)
txt(s, Inches(0.9), Inches(1.95), Inches(5.3), Inches(0.5), [("主指標（滲透率 + 學習深度）", 16, CYAN, True)])
bullets(s, Inches(0.9), Inches(2.5), Inches(5.3), Inches(1.5),
    ["AI 助教滲透率 11.7% → 目標翻倍", "AI 助教 DAU/MAU <10% → 20%+", "影片觀看次 / Session / 活躍天數"], size=14, gap=8)
card(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(2.3), CARD, line=ORANGE)
txt(s, Inches(7.1), Inches(1.95), Inches(5.3), Inches(0.5), [("A/B 驗證：遊戲化能否帶動看課？", 16, ORANGE, True)])
bullets(s, Inches(7.1), Inches(2.5), Inches(5.3), Inches(1.5),
    ["A 組：只玩知識挑戰（死路）", "B 組：EXP 綁定看課的養成", "比挑戰後 7 天回訪看課率 / 活躍天數 / 滲透率", "成功判準：B 顯著 > A"], size=13, gap=6)
card(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.1), NAVY2, line=GREEN)
txt(s, Inches(0.9), Inches(4.45), Inches(11.5), Inches(0.5), [("🔬 資料誠信聲明（技術評審加分項）", 16, GREEN, True)])
txt(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(1.4),
    [("曾觀察『有用 AI 者看課 1.83x』，但檢驗後中位數僅 1.10x、相關係數 r≈0.04，屬離群值假象。", 15, WHITE, False),
     ("→ 我們不主張此因果，改由 A/B 嚴謹驗證。寧可少講一個漂亮但脆弱的數字，", 15, GREY, False),
     ("   也不讓計畫在評審 Q&A 被一擊即破。", 15, GREY, False)])
footer(s, _idx[0])

# ---------------- 14. AWS 技術架構 ----------------
s = new()
header(s, "ARCHITECTURE", "AWS 技術架構", "13")
techs = [
    ("Amazon Bedrock", "RAG 撈 45 秒精華錨點 + 主動生成『知識缺口診斷 / 下一步路徑』", CYAN),
    ("Amazon DynamoDB", "去識別化 Session Token 儲存行為足跡矩陣，支撐長期記憶與去重", GREEN),
    ("向量檢索排除機制", "RAG 撈最匹配 3 堂，排除已看 >80% 的內容 ID（Filter Expression）", ORANGE),
    ("Amazon Titan (IP Indemnity)", "具智財賠償保證的繪圖模型 + 參數化紙娃娃系統，寵物生成 100% 商業安全", CYAN),
    ("隱私合規", "去識別化 + Opt-in 授權，符合 GDPR / CCPA", GREEN),
]
y = Inches(1.85)
for t, d, col in techs:
    card(s, Inches(0.6), y, Inches(12.1), Inches(0.9), CARD)
    rect(s, Inches(0.6), y, Inches(0.14), Inches(0.9), col)
    txt(s, Inches(0.95), y, Inches(3.6), Inches(0.9), [(t, 16, col, True)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.7), y, Inches(7.8), Inches(0.9), [(d, 14, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.0)
footer(s, _idx[0])

# ---------------- 15. Demo ----------------
s = new()
header(s, "DEMO", "可互動 Web/App 介面（聚焦完成度 15%）", "14")
txt(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.6),
    [("手機版 Web 原型（React + Tailwind），一條龍走完五大模組。", 17, CYAN, True)])
flow2 = ["知識PK擂台", "挑戰答題", "45秒軟著陸", "AI主動診斷", "推薦去重", "精靈養成", "多模態", "成效儀表板"]
x = Inches(0.6); y = Inches(2.7); w = Inches(2.9); h = Inches(0.85)
for i, f in enumerate(flow2):
    col = i % 4
    row = i // 4
    cx = Inches(0.6) + col * Inches(3.05)
    cy = Inches(2.7) + row * Inches(1.05)
    card(s, cx, cy, w, h, CARD, line=CYAN)
    txt(s, cx, cy, w, h, [(f"{i+1}. {f}", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
card(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.2), NAVY2, line=ORANGE)
txt(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(1.0),
    [("建議動線：擂台 → 答題(故意答錯) → 45秒軟著陸 → AI 主動診斷 → 推薦去重 → 精靈養成 → 多模態 → 成效", 14, WHITE, False),
     ("GitHub： github.com/a96020183/ppa-ai-hackathon-fish", 13, CYAN, True)])
footer(s, _idx[0])

# ---------------- 16. 自評評分表 ----------------
s = new()
header(s, "SELF-ASSESSMENT", "我們如何命中每一項評分", "15")
score = [
    ("商業影響力", "30%", "主指標用滲透率/DAU/學習深度；會員擴散＝擴大 AI 觸及池", ORANGE),
    ("創新度", "25%", "AI 被動→主動診斷、跨場景多模態，跳脫右側對話框", CYAN),
    ("資料洞察", "20%", "用①②⑥對話/軌跡資料證明斷點，資料驅動", GREEN),
    ("完成度", "15%", "Demo 聚焦一條龍主線，示意其餘模組", CYAN),
    ("技術性", "10%", "Bedrock RAG + 主動生成 + DynamoDB 去重 + Titan 合規生圖", ORANGE),
]
y = Inches(1.85)
for name, pct, desc, col in score:
    card(s, Inches(0.6), y, Inches(12.1), Inches(0.92), CARD)
    txt(s, Inches(0.95), y, Inches(2.6), Inches(0.92), [(name, 17, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.4), y, Inches(1.4), Inches(0.92), [(pct, 24, col, True)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.0), y, Inches(7.5), Inches(0.92), [(desc, 13, GREY, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.0)
footer(s, _idx[0])

# ---------------- 17. 結尾 ----------------
s = new(NAVY)
rect(s, 0, Inches(3.3), SW, Inches(0.1), ORANGE)
txt(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(1.0),
    [("讓 AI 更有用 → 驅動更多人看課 → 提升滲透與學習深度", 30, WHITE, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(0.8),
    [("團隊 fish · PressPlay AI 智慧學習生態系", 20, CYAN, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(0.6),
    [("Thank you  ·  Q & A", 22, ORANGE, True)], align=PP_ALIGN.CENTER)
footer(s, _idx[0])

prs.save(OUT)
print("SAVED:", OUT)
print("slides:", len(prs.slides._sldIdLst))
